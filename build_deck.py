"""Generate the VBOGS project slide deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

ROOT = "/home/oakley/ub/advanced_robotics/VBOGS"
OUT = os.path.join(ROOT, "VBOGS_Project_Overview.pptx")
DRIVE = os.path.join(ROOT, "outputs/v1_0/2013_05_28_drive_0007_sync")

# ---- palette -----------------------------------------------------------------
NAVY   = RGBColor(0x10, 0x1B, 0x33)
SLATE  = RGBColor(0x1E, 0x2A, 0x44)
ACCENT = RGBColor(0x4F, 0xC3, 0xF7)   # cyan
GOLD   = RGBColor(0xFF, 0xC1, 0x07)
WHITE  = RGBColor(0xF5, 0xF7, 0xFA)
GREY   = RGBColor(0xB8, 0xC2, 0xD0)
GREEN  = RGBColor(0x66, 0xBB, 0x6A)
CARD   = RGBColor(0x26, 0x33, 0x52)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def bg(slide, color=NAVY):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)
    return s


def box(slide, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold,italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        if isinstance(para, tuple):
            para = [para]
        for (t, sz, col, b, it) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = b; r.font.italic = it
            r.font.name = "Calibri"
    return tb


def R(t, sz=18, col=WHITE, b=False, it=False):
    return (t, sz, col, b, it)


def header(slide, kicker, title):
    box(slide, 0, 0, Inches(0.18), SH, fill=ACCENT)
    text(slide, Inches(0.55), Inches(0.30), Inches(11.5), Inches(0.4),
         [R(kicker, 14, ACCENT, True, False)])
    text(slide, Inches(0.55), Inches(0.58), Inches(12.2), Inches(0.9),
         [R(title, 32, WHITE, True, False)])
    box(slide, Inches(0.58), Inches(1.36), Inches(2.2), Inches(0.045), fill=GOLD)


def bullets(slide, x, y, w, h, items, size=17, gap=8):
    """items: list of (text, level) or text."""
    paras = []
    for it in items:
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        if lvl == 0:
            paras.append([R("▸  ", size, ACCENT, True), R(txt, size, WHITE)])
        elif lvl == 1:
            paras.append([R("        –  ", size-2, GOLD), R(txt, size-2, GREY)])
        else:  # heading
            paras.append([R(txt, size+1, GOLD, True)])
    text(slide, x, y, w, h, paras, space_after=gap, line_spacing=1.05)


def card(slide, x, y, w, h, title, body, tcol=ACCENT, fill=CARD):
    box(slide, x, y, w, h, fill=fill, line=SLATE, line_w=1.0, rounded=True)
    text(slide, x + Inches(0.18), y + Inches(0.12), w - Inches(0.36), Inches(0.5),
         [R(title, 16, tcol, True)])
    text(slide, x + Inches(0.18), y + Inches(0.52), w - Inches(0.36), h - Inches(0.66),
         [R(body, 12.5, GREY)], line_spacing=1.05)


def pagenum(slide, n):
    text(slide, SW - Inches(0.9), SH - Inches(0.5), Inches(0.6), Inches(0.3),
         [R(str(n), 12, GREY)], align=PP_ALIGN.RIGHT)


# =============================================================================
# Slide 1 — Title
# =============================================================================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
# accent band
box(s, 0, Inches(2.55), SW, Inches(0.06), fill=ACCENT)
box(s, 0, Inches(4.62), SW, Inches(0.06), fill=GOLD)
text(s, Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.5),
     [R("ADVANCED ROBOTICS  ·  PROJECT OVERVIEW", 16, ACCENT, True)])
text(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(1.2),
     [R("VBOGS", 76, WHITE, True)])
text(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(1.7),
     [[R("Variational Bayes Octree-GS", 34, WHITE, True)],
      [R("Uncertainty-aware 3D mapping & next-best-view planning", 24, GREY, False, True)]],
     space_after=6)
text(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(1.6),
     [[R("Octree-AnyGS scene scalability  +  Variational-Bayes Gaussian-Splatting uncertainty", 18, GREY)],
      [R("→  the next camera pose that maximally reduces what the map does not know", 18, GOLD, False, True)]],
     space_after=8)
text(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5),
     [R("Branch: Explicit-3D-Gaussians   ·   Dataset: KITTI-360 perspective stereo", 14, GREY, False, True)])

# =============================================================================
# Slide 2 — The Problem
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "MOTIVATION", "Why uncertainty-aware mapping?")
bullets(s, Inches(0.6), Inches(1.7), Inches(6.3), Inches(5.0), [
    ("The gap", 2),
    "Modern 3D Gaussian-Splatting reconstructs what the world looks like — but not where the model is guessing.",
    "An autonomous vehicle mapping a city needs to know which regions are poorly observed so it can go look again.",
    ("Two halves, never joined", 2),
    "Octree-GS gives scalable, level-of-detail scenes — but no notion of confidence.",
    "Variational-Bayes GS gives principled per-Gaussian uncertainty — but does not scale to driving-scale scenes.",
    ("VBOGS goal", 2),
    "Fuse both, then turn the uncertainty map into an actionable next-best-view (NBV) decision.",
], size=16, gap=9)
# right panel
box(s, Inches(7.25), Inches(1.7), Inches(5.45), Inches(4.9), fill=SLATE, rounded=True)
text(s, Inches(7.5), Inches(1.9), Inches(5.0), Inches(0.5),
     [R("The question VBOGS answers", 17, ACCENT, True)])
text(s, Inches(7.5), Inches(2.45), Inches(5.0), Inches(1.0),
     [R("“Given everything I have mapped so far, where should the camera look next?”",
        20, WHITE, True, True)], line_spacing=1.1)
box(s, Inches(7.5), Inches(3.55), Inches(4.95), Inches(0.04), fill=GOLD)
bullets(s, Inches(7.5), Inches(3.75), Inches(5.0), Inches(2.6), [
    "Knows what is mapped well",
    "Knows what is mapped poorly",
    "Scores candidate poses by expected information gain",
    "Outputs a ranked list of next camera poses",
], size=15, gap=10)
pagenum(s, 2)

# =============================================================================
# Slide 3 — Key Idea / Three pillars
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "CORE IDEA", "Three components, one pipeline")
cw, ch = Inches(3.95), Inches(3.5)
cy = Inches(1.85)
card(s, Inches(0.6), cy, cw, ch, "1 · Octree-AnyGS",
     "Scalable Gaussian-Splatting scene representation. Stores the scene as anchors — "
     "voxels at multiple levels of detail (LOD). Cell size shrinks by a fork factor per "
     "level: cur_size(l) = voxel_size / fork^l.\n\nPyTorch, SGD-trained on posed RGB "
     "frames. A single drive → ~243k anchors across 9 LOD levels.", tcol=ACCENT)
card(s, Inches(4.7), cy, cw, ch, "2 · VBGS uncertainty head",
     "Variational-Bayes Gaussian-Mixture model fit per anchor, post-hoc. Produces a full "
     "posterior — Normal-Wishart over (mean, covariance), Dirichlet over mixture weights "
     "— not just a point estimate.\n\nJAX. Posterior entropy → one uncertainty scalar "
     "per anchor.", tcol=GOLD)
card(s, Inches(8.8), cy, cw, ch, "3 · Next-best-view",
     "Per-anchor uncertainty is splatted through Octree-AnyGS's own LOD renderer into a "
     "pixel-space uncertainty image for each candidate camera pose.\n\nAlpha-normalised "
     "integration → a score per pose. Output: the ranked next-best views.", tcol=GREEN)
text(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.2),
     [[R("Design principle:  ", 16, GOLD, True), R(
        "the uncertainty head is bolted on after training — Octree-AnyGS Stage 1 is unchanged upstream. "
        "Stereo is used to build a depth cloud for the uncertainty fit, not as a training signal.", 16, GREY)]],
     line_spacing=1.15)
pagenum(s, 3)

# =============================================================================
# Slide 4 — Pipeline overview
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "ARCHITECTURE", "The five-stage pipeline")
stages = [
    ("STAGE 1", "Train scene", "Octree-AnyGS on\nposed RGB frames", ACCENT),
    ("STAGE 2", "Stereo cloud", "Disparity → 6D\nxyz+rgb point cloud", ACCENT),
    ("STAGE 3", "VBGS fit", "Per-anchor variational\nBayes GMM posterior", GOLD),
    ("STAGE 4", "Uncertainty", "Posterior entropy →\nscalar U[i] per anchor", GOLD),
    ("STAGE 5", "Next-best-view", "Splat U, score &\nrank candidate poses", GREEN),
]
x = Inches(0.55); y = Inches(2.0); bw = Inches(2.28); bh = Inches(2.5); gap = Inches(0.18)
for i, (kick, title, body, col) in enumerate(stages):
    bx = x + (bw + gap) * i
    box(s, bx, y, bw, bh, fill=CARD, line=col, line_w=1.5, rounded=True)
    box(s, bx, y, bw, Inches(0.5), fill=col, rounded=True)
    text(s, bx, y + Inches(0.04), bw, Inches(0.42),
         [R(kick, 13, NAVY, True)], align=PP_ALIGN.CENTER)
    text(s, bx + Inches(0.1), y + Inches(0.62), bw - Inches(0.2), Inches(0.5),
         [R(title, 18, WHITE, True)], align=PP_ALIGN.CENTER)
    text(s, bx + Inches(0.12), y + Inches(1.15), bw - Inches(0.24), Inches(1.2),
         [R(body, 13, GREY)], align=PP_ALIGN.CENTER, line_spacing=1.1)
    if i < 4:
        text(s, bx + bw - Inches(0.02), y + Inches(0.85), gap + Inches(0.2), Inches(0.5),
             [R("▶", 18, GOLD, True)], align=PP_ALIGN.CENTER)
# CLI strip
box(s, Inches(0.55), Inches(4.85), Inches(12.2), Inches(1.45), fill=SLATE, rounded=True)
text(s, Inches(0.8), Inches(4.98), Inches(11.6), Inches(0.4),
     [R("Orchestrated end-to-end by  scripts/run_drive_pipeline.py", 16, ACCENT, True)])
text(s, Inches(0.8), Inches(5.38), Inches(11.8), Inches(0.9),
     [[R("prepare → train → stereo → bucket → fit → inspect → uncertainty → map-viz → render → nbv → bundle",
         14, WHITE, False, False)],
      [R("Resumable at any stage (--start-at / --stop-after); dispatches each stage to the right Docker service & GPU.",
         13, GREY, False, True)]], space_after=5)
pagenum(s, 4)

# =============================================================================
# Slide 5 — Stage 1: Train scene
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "STAGE 1  ·  TRAIN THE SCENE", "Octree-AnyGS scene representation")
bullets(s, Inches(0.6), Inches(1.7), Inches(6.6), Inches(5.0), [
    ("Input", 2),
    "Posed RGB frames from KITTI-360 — KITTI-360 layout is adapted to COLMAP format by prepare_kitti360_colmap.py.",
    ("Representation", 2),
    "Scene = a flat tensor of anchors; each anchor carries a level l. The anchor IS the voxel at its LOD — no separate voxel object.",
    "cur_size(l) = voxel_size / fork^l  — finer levels = smaller cells.",
    "Each anchor spawns n_offsets Gaussians at render time.",
    ("Training", 2),
    "Standard Octree-AnyGS SGD training, unchanged upstream. Progressive LOD growth.",
    "Two appearance back-ends: implicit3D (neural MLP) vs explicit3D (explicit SH) — see slide 12.",
], size=15, gap=7)
# stats panel
box(s, Inches(7.5), Inches(1.7), Inches(5.2), Inches(4.7), fill=SLATE, rounded=True)
text(s, Inches(7.75), Inches(1.88), Inches(4.7), Inches(0.4),
     [R("Example — drive 2013_05_28_drive_0007", 15, ACCENT, True)])
stat_rows = [
    ("Frames trained", "1,000"),
    ("Total anchors", "243,377"),
    ("LOD levels", "9"),
    ("Gaussians / anchor", "n_offsets"),
    ("Optimizer", "SGD (PyTorch)"),
]
ry = Inches(2.4)
for k, v in stat_rows:
    text(s, Inches(7.75), ry, Inches(3.0), Inches(0.45), [R(k, 15, GREY)])
    text(s, Inches(10.6), ry, Inches(1.9), Inches(0.45),
         [R(v, 16, GOLD, True)], align=PP_ALIGN.RIGHT)
    box(s, Inches(7.75), ry + Inches(0.42), Inches(4.7), Emu(9525), fill=CARD)
    ry += Inches(0.66)
text(s, Inches(7.75), ry + Inches(0.05), Inches(4.7), Inches(0.8),
     [R("Octree-AnyGS is included as a read-only submodule; VBOGS only reuses its grid & geometry path.",
        12.5, GREY, False, True)], line_spacing=1.1)
pagenum(s, 5)

# =============================================================================
# Slide 6 — Stage 2: Stereo point cloud
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "STAGE 2  ·  STEREO POINT CLOUD", "From disparity to a 6D RGB point cloud")
bullets(s, Inches(0.6), Inches(1.65), Inches(6.7), Inches(3.4), [
    ("Per stereo pair", 2),
    "Disparity = StereoMatch(img_L, img_R) — pluggable matcher (default OpenCV SGBM; RAFT-Stereo reserved).",
    "depth_L = (K_L.fx · baseline) / disparity",
    "Unproject to camera frame, then transform to world coords with the camera pose.",
    "Concatenate xyz + rgb → a 6D point.",
    ("Quality filters", 2),
    "Left-right consistency check, texture threshold, min-disparity & max-depth gates drop low-confidence pixels.",
], size=15, gap=7)
# formula card
box(s, Inches(0.6), Inches(5.05), Inches(6.7), Inches(1.5), fill=CARD, line=ACCENT, line_w=1.2, rounded=True)
text(s, Inches(0.85), Inches(5.18), Inches(6.2), Inches(0.4),
     [R("Two coordinate systems — kept in parallel", 14, ACCENT, True)])
text(s, Inches(0.85), Inches(5.55), Inches(6.2), Inches(0.95),
     [[R("points_world", 14, GOLD, True), R("  — bucket points into the Octree grid", 13, GREY)],
      [R("points_norm", 14, GOLD, True), R("   — normalized scale that VBGS priors expect", 13, GREY)]],
     space_after=4)
# stats panel
box(s, Inches(7.6), Inches(1.65), Inches(5.1), Inches(4.9), fill=SLATE, rounded=True)
text(s, Inches(7.85), Inches(1.82), Inches(4.6), Inches(0.4),
     [R("Drive 0007 — stereo stats", 15, ACCENT, True)])
rows = [
    ("Matcher", "SGBM"),
    ("Frames", "1,000"),
    ("Points produced", "100,000,000"),
    ("Image size", "1408 × 376"),
    ("Baseline", "0.594 m"),
    ("fx", "552.55 px"),
    ("Max depth", "80 m"),
    ("LR-consistency", "1.0 px"),
]
ry = Inches(2.32)
for k, v in rows:
    text(s, Inches(7.85), ry, Inches(2.7), Inches(0.4), [R(k, 14, GREY)])
    text(s, Inches(10.4), ry, Inches(2.0), Inches(0.4),
         [R(v, 14.5, GOLD, True)], align=PP_ALIGN.RIGHT)
    box(s, Inches(7.85), ry + Inches(0.38), Inches(4.6), Emu(9525), fill=CARD)
    ry += Inches(0.52)
pagenum(s, 6)

# =============================================================================
# Slide 7 — Stage 3: Bucketing + VBGS fit
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "STAGE 3  ·  PER-ANCHOR VBGS FIT", "Bucket the cloud, fit a posterior per anchor")
bullets(s, Inches(0.6), Inches(1.6), Inches(6.5), Inches(2.7), [
    ("Bucketing", 2),
    "Each point is assigned to every anchor whose cell contains it — at all LOD levels, not just the finest.",
    "Same integer-grid math Octree-AnyGS uses, so buckets match anchors exactly.",
    "Why all levels: a coarse anchor must reflect what it would see if its LOD were chosen — finest-only would starve coarse anchors and make them look spuriously uncertain.",
], size=14.5, gap=6)
bullets(s, Inches(0.6), Inches(4.35), Inches(6.5), Inches(2.8), [
    ("Variational-Bayes GMM per anchor", 2),
    "Anchors with < MIN_POINTS_PER_ANCHOR (20) → flagged unobserved.",
    "VBEM fit in batches via vbgs.fit_gmm_step (streaming sufficient statistics).",
    "K-growth: start K = K_INIT, double K while per-point ELBO gain ≥ tolerance, cap at K_MAX.",
], size=14.5, gap=6)
# right: K-growth diagram
box(s, Inches(7.35), Inches(1.6), Inches(5.35), Inches(5.0), fill=SLATE, rounded=True)
text(s, Inches(7.6), Inches(1.75), Inches(4.9), Inches(0.4),
     [R("Adaptive model complexity (K-growth)", 15, ACCENT, True)])
ladder = [("K = 10", "start (K_INIT)"), ("K = 20", "ΔELBO ≥ tol → grow"),
          ("K = 40", "cap (K_MAX)")]
ly = Inches(2.3)
for kk, note in ladder:
    box(s, Inches(7.7), ly, Inches(1.55), Inches(0.55), fill=CARD, line=GOLD, line_w=1.2, rounded=True)
    text(s, Inches(7.7), ly + Inches(0.06), Inches(1.55), Inches(0.45),
         [R(kk, 16, GOLD, True)], align=PP_ALIGN.CENTER)
    text(s, Inches(9.45), ly + Inches(0.08), Inches(3.0), Inches(0.45),
         [R(note, 13, GREY)])
    if note != ladder[-1][1]:
        text(s, Inches(8.3), ly + Inches(0.5), Inches(0.4), Inches(0.35),
             [R("▼", 13, ACCENT)], align=PP_ALIGN.CENTER)
    ly += Inches(0.95)
box(s, Inches(7.6), Inches(5.1), Inches(4.85), Inches(0.035), fill=GOLD)
text(s, Inches(7.6), Inches(5.25), Inches(4.9), Inches(1.3),
     [[R("Posterior per component:  ", 13.5, ACCENT, True),
       R("Normal-Wishart(mean, κ, U, n) over (μ, Σ); Dirichlet α over mixture weights π.",
         13.5, GREY)]], line_spacing=1.15)
text(s, Inches(7.6), Inches(6.05), Inches(4.9), Inches(0.45),
     [R("Framework bridge: PyTorch ⇄ JAX via numpy arrays on disk — no shared autograd.",
        12, GREY, False, True)])
pagenum(s, 7)

# =============================================================================
# Slide 8 — Stage 4: Uncertainty scalar
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "STAGE 4  ·  REDUCE TO A SCALAR", "Posterior entropy → one number per anchor")
bullets(s, Inches(0.6), Inches(1.65), Inches(6.5), Inches(3.6), [
    ("Observed anchors", 2),
    "U[i] = Σ_k  π̄[k] · H_k   — expected posterior entropy, mixture-weighted.",
    "H_k = closed-form Normal-Wishart entropy + delta (colour) entropy of component k.",
    "π̄ = E[π] from the Dirichlet posterior over mixture weights.",
    ("Unobserved anchors", 2),
    "U[i] = U_MAX — anchors with too few points are treated as maximally uncertain.",
    ("Why it is comparable", 2),
    "All entropies computed in normalized coords → a shared unit scale across every anchor.",
], size=15, gap=7)
# interpretation card
box(s, Inches(0.6), Inches(5.45), Inches(6.5), Inches(1.1), fill=CARD, line=GOLD, line_w=1.2, rounded=True)
text(s, Inches(0.82), Inches(5.55), Inches(6.05), Inches(0.95),
     [[R("Intuition:  ", 14, GOLD, True),
       R("high κ & n → confident location & shape → low entropy → low U. "
         "Sparse / never-seen anchors → high U.", 13.5, GREY)]], line_spacing=1.15)
# embed histogram
hist = os.path.join(DRIVE, "uncertainty/uncertainty_histogram.png")
box(s, Inches(7.35), Inches(1.65), Inches(5.4), Inches(4.95), fill=SLATE, rounded=True)
text(s, Inches(7.6), Inches(1.78), Inches(4.9), Inches(0.4),
     [R("Drive 0007 — uncertainty distribution", 15, ACCENT, True)])
if os.path.exists(hist):
    s.shapes.add_picture(hist, Inches(7.6), Inches(2.25), width=Inches(4.9))
text(s, Inches(7.6), Inches(5.05), Inches(4.9), Inches(1.45),
     [[R("243,377", 15, GOLD, True), R(" anchors total   ", 13, GREY),
       R("·  94,473", 15, GOLD, True), R(" observed", 13, GREY)],
      [R("148,904", 15, GOLD, True), R(" unobserved → U_MAX ≈ 37.17", 13, GREY)],
      [R("observed U: p50 = 23.3, p90 = 31.7, p98 = 34.4", 12.5, GREY, False, True)]],
     space_after=4, line_spacing=1.1)
pagenum(s, 8)

# =============================================================================
# Slide 9 — Stage 5: NBV selection
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "STAGE 5  ·  NEXT-BEST-VIEW", "Splat the uncertainty, score the poses")
bullets(s, Inches(0.6), Inches(1.6), Inches(6.4), Inches(3.4), [
    ("render_scalar()", 2),
    "Reuses Octree-AnyGS LOD mask + generate_gaussians geometry path — but substitutes the per-anchor uncertainty scalar in place of MLP colour.",
    "Per-anchor U is broadcast over n_offsets, masked, and rasterized with gsplat.",
    "Returns an uncertainty image + an alpha image per candidate pose. Occlusion & LOD handled for free.",
    ("Scoring each candidate", 2),
    "score = Σ uncertainty_image / (Σ alpha_image + ε)",
    "Alpha-normalisation rewards poses that see uncertain surfaces — not poses that merely see a lot of surface.",
], size=14.5, gap=6)
# formula highlight
box(s, Inches(0.6), Inches(5.35), Inches(6.4), Inches(1.15), fill=CARD, line=GREEN, line_w=1.4, rounded=True)
text(s, Inches(0.6), Inches(5.5), Inches(6.4), Inches(0.5),
     [R("score  =  Σ U_image  /  (Σ α_image + ε)", 22, GREEN, True)], align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(5.98), Inches(6.4), Inches(0.4),
     [R("expected posterior entropy per unit of visible surface", 13, GREY, False, True)],
     align=PP_ALIGN.CENTER)
# embed nbv heat + alpha
heat = os.path.join(DRIVE, "nbv/viz/rank_01_unc_heat.png")
alpha = os.path.join(DRIVE, "nbv/viz/rank_01_alpha.png")
box(s, Inches(7.25), Inches(1.6), Inches(5.5), Inches(4.95), fill=SLATE, rounded=True)
text(s, Inches(7.5), Inches(1.73), Inches(5.0), Inches(0.4),
     [R("Top-ranked candidate — rank 01", 15, ACCENT, True)])
if os.path.exists(heat):
    s.shapes.add_picture(heat, Inches(7.5), Inches(2.2), width=Inches(5.0))
if os.path.exists(alpha):
    s.shapes.add_picture(alpha, Inches(7.5), Inches(4.05), width=Inches(5.0))
text(s, Inches(7.5), Inches(5.9), Inches(5.0), Inches(0.6),
     [R("uncertainty image (top)  ·  alpha / opacity image (bottom)", 12.5, GREY, False, True)],
     align=PP_ALIGN.CENTER)
pagenum(s, 9)

# =============================================================================
# Slide 10 — NBV results
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "RESULTS  ·  NEXT-BEST-VIEW", "Ranked candidates for drive 0007")
# left: leaderboard
box(s, Inches(0.6), Inches(1.7), Inches(6.3), Inches(4.9), fill=SLATE, rounded=True)
text(s, Inches(0.85), Inches(1.85), Inches(5.8), Inches(0.4),
     [R("Top-5 of 125 test-camera candidates", 16, ACCENT, True)])
# header row
hy = Inches(2.35)
cols = [("RANK", Inches(0.85), Inches(0.9)), ("IMAGE", Inches(1.85), Inches(1.6)),
        ("SCORE", Inches(3.55), Inches(1.3)), ("VIS. ANCHORS", Inches(4.95), Inches(1.6))]
for label, cx, cwid in cols:
    text(s, cx, hy, cwid, Inches(0.35), [R(label, 12, GOLD, True)])
box(s, Inches(0.85), Inches(2.7), Inches(5.8), Emu(12700), fill=GOLD)
nbv_rows = [
    ("1", "0000000609", "28.63", "41,414"),
    ("2", "0000000617", "28.46", "41,506"),
    ("3", "0000000601", "27.9", "—"),
    ("4", "0000000625", "27.6", "—"),
    ("5", "0000000593", "27.3", "—"),
]
ry = Inches(2.85)
for rank, img, score, anc in nbv_rows:
    col = GREEN if rank == "1" else WHITE
    text(s, Inches(0.85), ry, Inches(0.9), Inches(0.4), [R("#" + rank, 15, col, True)])
    text(s, Inches(1.85), ry, Inches(1.7), Inches(0.4), [R(img, 14, GREY)])
    text(s, Inches(3.55), ry, Inches(1.3), Inches(0.4), [R(score, 15, col, rank == "1")])
    text(s, Inches(4.95), ry, Inches(1.6), Inches(0.4), [R(anc, 14, GREY)])
    box(s, Inches(0.85), ry + Inches(0.42), Inches(5.8), Emu(9525), fill=CARD)
    ry += Inches(0.6)
text(s, Inches(0.85), Inches(6.0), Inches(5.8), Inches(0.5),
     [R("Ranks 3–5 scores approximate; full list in nbv/nbv_scores.json", 11.5, GREY, False, True)])
# right: winner card
box(s, Inches(7.25), Inches(1.7), Inches(5.45), Inches(4.9), fill=CARD, line=GREEN, line_w=1.6, rounded=True)
text(s, Inches(7.5), Inches(1.88), Inches(5.0), Inches(0.4),
     [R("★  SELECTED NEXT-BEST VIEW", 16, GREEN, True)])
text(s, Inches(7.5), Inches(2.3), Inches(5.0), Inches(0.5),
     [R("Frame 0000000609  ·  candidate #76", 18, WHITE, True)])
win = [
    ("Score", "28.63"),
    ("Σ uncertainty", "770,906"),
    ("Σ alpha", "26,929"),
    ("Visible anchors", "41,414"),
    ("Visible Gaussians", "40,504"),
    ("Camera centre", "(361.2, −2580.1, 139.9)"),
]
wy = Inches(2.95)
for k, v in win:
    text(s, Inches(7.55), wy, Inches(2.4), Inches(0.4), [R(k, 14, GREY)])
    text(s, Inches(9.95), wy, Inches(2.55), Inches(0.4),
         [R(v, 14.5, GOLD, True)], align=PP_ALIGN.RIGHT)
    box(s, Inches(7.55), wy + Inches(0.4), Inches(4.9), Emu(9525), fill=SLATE)
    wy += Inches(0.56)
text(s, Inches(7.55), wy + Inches(0.05), Inches(4.9), Inches(0.6),
     [R("This pose maximises expected posterior-entropy reduction per unit of visible surface.",
        12.5, GREY, False, True)], line_spacing=1.1)
pagenum(s, 10)

# =============================================================================
# Slide 11 — Rendered uncertainty views
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "RESULTS  ·  DIAGNOSTIC VIEWS", "Rendered uncertainty across the drive")
text(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.6),
     [R("render_uncertainty_views.py produces side-by-side RGB ↔ uncertainty heatmaps for every "
        "camera — 1,000 train + 125 test views for drive 0007.", 15, GREY)], line_spacing=1.1)
sbs_dir = os.path.join(DRIVE, "views/test/side_by_side")
picks = ["00041_0000000329.png", "00073_0000000585.png", "00099_0000000793.png"]
yy = Inches(2.35)
for i, name in enumerate(picks):
    p = os.path.join(sbs_dir, name)
    px = Inches(0.6) + (Inches(4.1)) * i
    box(s, px, yy, Inches(3.95), Inches(1.55), fill=SLATE, rounded=True)
    if os.path.exists(p):
        s.shapes.add_picture(p, px + Inches(0.07), yy + Inches(0.32),
                             width=Inches(3.81))
    text(s, px, yy + Inches(0.04), Inches(3.95), Inches(0.3),
         [R("test view  " + name.split("_")[1].replace(".png", ""), 12, ACCENT, True)],
         align=PP_ALIGN.CENTER)
text(s, Inches(0.6), Inches(4.25), Inches(12.1), Inches(0.5),
     [R("Left = RGB render   ·   Right = uncertainty heatmap (warm = uncertain)", 13, GREY, False, True)])
# takeaways
box(s, Inches(0.6), Inches(4.85), Inches(12.1), Inches(1.7), fill=CARD, rounded=True)
text(s, Inches(0.85), Inches(4.98), Inches(11.6), Inches(0.4),
     [R("What the heatmaps show", 15, GOLD, True)])
bullets(s, Inches(0.85), Inches(5.35), Inches(11.6), Inches(1.1), [
    "Distant facades, sky boundaries and sparsely-observed side streets light up as high-uncertainty.",
    "Well-traversed road surface and near-field geometry stay cool — confirming the posterior-entropy signal is spatially meaningful.",
], size=14, gap=5)
pagenum(s, 11)

# =============================================================================
# Slide 12 — implicit vs explicit (this branch)
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "BRANCH FOCUS", "Explicit-3D-Gaussians: implicit3D vs explicit3D")
text(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.55),
     [R("The current branch adds a --gaussian-type switch and the configs to compare the two "
        "Octree-AnyGS appearance back-ends head-to-head.", 15, GREY)], line_spacing=1.1)
# two columns
cw = Inches(5.95)
card(s, Inches(0.6), Inches(2.3), cw, Inches(3.4), "implicit3D  (neural default)",
     "Anchors store feature vectors; opacity, covariance and colour come from small MLPs.\n\n"
     "• feat_dim = 16\n"
     "• MLP learning rates for opacity / colour / cov\n"
     "• More expressive — MLPs can absorb appearance artifacts\n"
     "• Higher memory & indirection", tcol=ACCENT)
card(s, Inches(6.78), Inches(2.3), cw, Inches(3.4), "explicit3D  (this branch)",
     "Anchors store explicit spherical-harmonic colour, opacity, scale and rotation directly.\n\n"
     "• feature_lr 0.0025 · opacity_lr 0.05\n"
     "• scaling_lr 0.005 · rotation_lr 0.001 · lambda_dreg 0\n"
     "• Direct parameterisation, clearer geometry\n"
     "• Potentially faster, lighter inference", tcol=GOLD)
box(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.95), fill=SLATE, rounded=True)
text(s, Inches(0.85), Inches(6.0), Inches(11.6), Inches(0.8),
     [[R("Experiment configs:  ", 14, ACCENT, True),
       R("exp02-A_explicit3d_baseline  vs  exp02_B_implicit3d_baseline  — identical pipeline, "
         "only the Gaussian type differs, so downstream uncertainty & NBV are directly comparable.",
         13.5, GREY)]], line_spacing=1.15)
pagenum(s, 12)

# =============================================================================
# Slide 13 — System / tech stack
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "SYSTEM", "Containerised, multi-framework stack")
text(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.5),
     [R("PyTorch (Octree-AnyGS) and JAX (vbgs) have conflicting CUDA builds — so each runs in its "
        "own service, coordinated by a third.", 15, GREY)], line_spacing=1.1)
svc = [
    ("vbogs-torch", "PyTorch · CUDA", "Octree-AnyGS training, stereo matching, point bucketing, scalar rendering.", ACCENT),
    ("vbogs-jax", "JAX · GPU", "VBGS per-anchor variational fitting, fit inspection, uncertainty scalar computation.", GOLD),
    ("vbogs-pipeline", "Orchestrator", "Runs run_drive_pipeline.py, dispatches each stage to the right service, bundles & exports outputs.", GREEN),
]
x = Inches(0.6); y = Inches(2.3); bw = Inches(3.95); bh = Inches(2.6); gap = Inches(0.18)
for i, (name, sub, body, col) in enumerate(svc):
    bx = x + (bw + gap) * i
    box(s, bx, y, bw, bh, fill=CARD, line=col, line_w=1.5, rounded=True)
    box(s, bx, y, bw, Inches(0.62), fill=col, rounded=True)
    text(s, bx + Inches(0.15), y + Inches(0.04), bw - Inches(0.3), Inches(0.35),
         [R(name, 17, NAVY, True)])
    text(s, bx + Inches(0.15), y + Inches(0.34), bw - Inches(0.3), Inches(0.3),
         [R(sub, 12, NAVY, True)])
    text(s, bx + Inches(0.15), y + Inches(0.78), bw - Inches(0.3), Inches(1.7),
         [R(body, 13.5, GREY)], line_spacing=1.15)
# bottom strip
box(s, Inches(0.6), Inches(5.2), Inches(12.13), Inches(1.5), fill=SLATE, rounded=True)
text(s, Inches(0.85), Inches(5.32), Inches(11.6), Inches(0.4),
     [R("Stack details", 15, GOLD, True)])
bullets(s, Inches(0.85), Inches(5.68), Inches(11.7), Inches(1.0), [
    "Shared external Docker volumes: KITTI-360 · COLMAP · OCTREE-ANYGS — created once, mounted in every service.",
    "Deployable via Docker Compose or Portainer (build-from-repository template); optional Google-Drive export of bundled outputs.",
], size=13.5, gap=4)
pagenum(s, 13)

# =============================================================================
# Slide 14 — Repo layout / artifacts
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "ENGINEERING", "Codebase & per-drive artifacts")
# left: code
box(s, Inches(0.6), Inches(1.7), Inches(6.05), Inches(4.9), fill=SLATE, rounded=True)
text(s, Inches(0.85), Inches(1.82), Inches(5.5), Inches(0.4),
     [R("Repository structure", 16, ACCENT, True)])
code_items = [
    ("scripts/", "11 stage entry-points + orchestrator", GOLD),
    ("  run_drive_pipeline.py", "master orchestration (M1–M7)", GREY),
    ("  stereo_to_pointcloud.py", "Stage 2", GREY),
    ("  bucket_points.py / fit_anchors.py", "Stage 3", GREY),
    ("  compute_uncertainty.py", "Stage 4", GREY),
    ("  score_nbv.py", "Stage 5", GREY),
    ("vbogs/", "shared library", GOLD),
    ("  render.py", "render_scalar — scalar splatting", GREY),
    ("  data_layout.py / io.py / fit_planning.py", "", GREY),
    ("Octree-AnyGS/  ·  vbgs/", "read-only submodules", GOLD),
    ("docs/manuscript/Algorithm.tex", "authoritative algorithm spec", GOLD),
    ("docker/  ·  *compose*.yml", "containerisation", GOLD),
]
cy = Inches(2.3)
for name, desc, col in code_items:
    text(s, Inches(0.9), cy, Inches(3.4), Inches(0.35),
         [R(name, 13, col, col != GREY)])
    text(s, Inches(4.0), cy, Inches(2.5), Inches(0.35),
         [R(desc, 11.5, GREY, False, True)])
    cy += Inches(0.345)
# right: artifacts
box(s, Inches(6.85), Inches(1.7), Inches(5.85), Inches(4.9), fill=SLATE, rounded=True)
text(s, Inches(7.1), Inches(1.82), Inches(5.3), Inches(0.4),
     [R("Per-drive output bundle", 16, GOLD, True)])
arti = [
    "pointclouds/stereo/ — 6D world point cloud (NPZ + PLY)",
    "pointclouds/anchors/ — uncertainty-coloured anchors",
    "uncertainty/ — U.npy, components, histogram",
    "views/ — RGB ↔ uncertainty side-by-sides (train+test)",
    "nbv/ — nbv_scores.json, top images, viz heatmaps",
    "prepared/ + octree/ — COLMAP metadata & train config",
    "run_manifest.json — full artifact provenance",
]
bullets(s, Inches(7.1), Inches(2.35), Inches(5.4), Inches(3.0), arti, size=13.5, gap=8)
box(s, Inches(7.1), Inches(5.35), Inches(5.35), Inches(0.035), fill=GOLD)
text(s, Inches(7.1), Inches(5.5), Inches(5.4), Inches(1.0),
     [R("bundle_run_outputs.py curates everything into a versioned, zipped bundle with a manifest — "
        "heavy checkpoints stay in their native volumes, referenced by path.", 13, GREY, False, True)],
     line_spacing=1.15)
pagenum(s, 14)

# =============================================================================
# Slide 15 — Limitations & future
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "DISCUSSION", "Known limitations & future work")
# limitations
box(s, Inches(0.6), Inches(1.75), Inches(5.95), Inches(4.8), fill=CARD, line=GOLD, line_w=1.2, rounded=True)
text(s, Inches(0.85), Inches(1.9), Inches(5.5), Inches(0.4),
     [R("⚠  Limitations (acknowledged in design)", 16, GOLD, True)])
bullets(s, Inches(0.85), Inches(2.4), Inches(5.5), Inches(4.0), [
    "Empty-region blindness — render_scalar only splats through existing anchors. Truly never-seen volumes score zero, so NBV steers toward poorly-modelled regions, not unknown space.",
    "ELBO-biased K-selection — the KL term scales with K, so per-point mean ELBO is a pragmatic, not information-theoretically clean, model-selection criterion.",
    "Normalization dependence — cross-anchor entropy comparison only holds because everything is fit in shared normalized coords.",
], size=13.5, gap=10)
# future
box(s, Inches(6.75), Inches(1.75), Inches(5.95), Inches(4.8), fill=CARD, line=GREEN, line_w=1.2, rounded=True)
text(s, Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.4),
     [R("→  Future directions", 16, GREEN, True)])
bullets(s, Inches(7.0), Inches(2.4), Inches(5.5), Inches(4.0), [
    "Volumetric occupancy prior or per-pixel “unknown ray” penalty to enable true exploration into empty space.",
    "Held-out log-likelihood or BIC for principled K-selection.",
    "RAFT-Stereo matcher for denser, longer-range depth.",
    "Continual variant — fit_gmm_step already supports streaming updates as new frames arrive.",
    "Complete M7: human qualitative validation of NBV picks.",
], size=13.5, gap=9)
pagenum(s, 15)

# =============================================================================
# Slide 16 — Summary
# =============================================================================
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0, Inches(2.45), SW, Inches(0.06), fill=ACCENT)
text(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.5),
     [R("SUMMARY", 16, ACCENT, True)])
text(s, Inches(0.9), Inches(1.1), Inches(11.5), Inches(1.3),
     [R("VBOGS in one slide", 44, WHITE, True)])
takeaways = [
    ("What it is", "A post-hoc Bayesian uncertainty head on a scalable Octree-GS scene, turned into a next-best-view planner for autonomous-vehicle mapping."),
    ("How it works", "5 stages: train scene → stereo cloud → per-anchor variational-Bayes GMM → posterior-entropy scalar → splat & score candidate poses."),
    ("Why it matters", "The map doesn't just know what the world looks like — it knows where it is uncertain, and can act on it."),
    ("Status", "M1–M6 implemented & containerised; demonstrated end-to-end on KITTI-360 drive 0007 (243k anchors, 100M stereo points, 125 NBV candidates ranked)."),
]
ty = Inches(2.85)
for head, body in takeaways:
    box(s, Inches(0.9), ty, Inches(2.55), Inches(0.78), fill=CARD, line=GOLD, line_w=1.0, rounded=True)
    text(s, Inches(0.9), ty + Inches(0.16), Inches(2.55), Inches(0.5),
         [R(head, 16, GOLD, True)], align=PP_ALIGN.CENTER)
    text(s, Inches(3.7), ty + Inches(0.02), Inches(8.7), Inches(0.85),
         [R(body, 15.5, WHITE)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    ty += Inches(0.95)
text(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4),
     [R("VBOGS  ·  Variational Bayes Octree-GS  ·  Branch: Explicit-3D-Gaussians", 13, GREY, False, True)])
pagenum(s, 16)

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
